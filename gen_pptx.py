# gen_pptx.py
from typing import List, Dict, Any, Optional, Tuple
import io
import os
import tempfile

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Optional URL support for Image: if 'requests' isn't installed, URLs are ignored safely.
try:
    import requests  # type: ignore
    _HAS_REQUESTS = True
except Exception:
    _HAS_REQUESTS = False


# -------------------------
# Basic helpers
# -------------------------

def _to_lines(val) -> List[str]:
    """Normalize strings/lists to a list of non-empty, stripped lines."""
    if not val:
        return []
    if isinstance(val, str):
        return [ln.strip() for ln in val.splitlines() if ln.strip()]
    if isinstance(val, list):
        out: List[str] = []
        for v in val:
            if isinstance(v, str):
                out.extend([ln.strip() for ln in v.splitlines() if ln.strip()])
        return out
    return []


def _clear_text_frame(tf):
    """Clear a text frame, keep a clean start, and disable autofit."""
    if tf is None:
        return
    try:
        tf.clear()
    except Exception:
        # best-effort fallback
        try:
            if tf.paragraphs:
                tf.paragraphs[0].text = ""
                for p in tf.paragraphs[1:]:
                    try:
                        tf._element.remove(p._p)
                    except Exception:
                        pass
            else:
                tf.add_paragraph()
        except Exception:
            pass
    try:
        tf.word_wrap = True
    except Exception:
        pass
    # Disable auto-size if available (prevents inconsistent font shrinking)
    try:
        tf.auto_size = False  # type: ignore[attr-defined]
    except Exception:
        pass


def _hex_to_rgb(hexstr: str) -> RGBColor:
    s = (hexstr or "").strip().lstrip("#")
    if len(s) == 3:
        s = "".join(ch * 2 for ch in s)
    try:
        r = int(s[0:2], 16); g = int(s[2:4], 16); b = int(s[4:6], 16)
        return RGBColor(r, g, b)
    except Exception:
        # default dark gray if parsing fails
        return RGBColor(17, 17, 17)


def _apply_font(p, *, size_pt: int, bold: bool, name: str, color: RGBColor, bullet: bool = False, level: int = 0):
    """Apply consistent font + bullet state to a paragraph."""
    if p is None:
        return
    p.font.size = Pt(size_pt)
    p.font.bold = bold
    try:
        p.font.name = name
    except Exception:
        pass
    try:
        p.font.color.rgb = color
    except Exception:
        pass
    # Bullet state
    try:
        p.level = level
        p._set_bullet(bool(bullet))  # supported across python-pptx versions
    except Exception:
        if bullet and not getattr(p, "text", "").startswith("• "):
            p.text = f"• {p.text}"


def _add_para(tf, text: str, *, size_pt: int, bold: bool, name: str, color: RGBColor):
    """Append a normal paragraph (bullets OFF)."""
    if tf is None:
        return
    p = tf.add_paragraph()
    p.text = text
    _apply_font(p, size_pt=size_pt, bold=bold, name=name, color=color, bullet=False, level=0)


def _add_bullet(tf, text: str, *, size_pt: int, level: int, name: str, color: RGBColor):
    """Append a bulleted paragraph."""
    if tf is None:
        return
    p = tf.add_paragraph()
    p.text = text
    _apply_font(p, size_pt=size_pt, bold=False, name=name, color=color, bullet=True, level=level)


def _ph_val(name: str):
    """Get placeholder enum value if it exists in this python-pptx version."""
    return getattr(PP_PLACEHOLDER, name, None)


def _find_title_placeholder(slide):
    """Find a title-like placeholder."""
    for nm in ("TITLE", "CENTER_TITLE"):
        tv = _ph_val(nm)
        if tv is None:
            continue
        for shp in getattr(slide, "placeholders", []):
            try:
                if shp.placeholder_format.type == tv:
                    return shp
            except Exception:
                pass
    return None


def _content_bounds(prs: Presentation) -> Tuple[int, int, int, int]:
    """Return left, top, width, height with safe margins relative to slide size."""
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    left_margin = Inches(0.75)
    right_margin = Inches(0.75)
    top_margin = Inches(1.8)   # leave room under title
    bottom_margin = Inches(1.0)
    left = left_margin
    top = top_margin
    width = slide_w - left_margin - right_margin
    height = slide_h - top_margin - bottom_margin
    return left, top, width, height


def _apply_slide_bg(slide, color_rgb: RGBColor):
    """Fill the whole slide background with a solid color."""
    try:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color_rgb
    except Exception:
        pass


# -------------------------
# Auto-split helper
# -------------------------

def _chunk_content(text_lines: List[str], bullet_lines: List[str], max_text=6, max_bullets=6):
    """
    Yield (text_chunk, bullets_chunk) enforcing per-slide limits.
    If both text and bullets exist, allocate capacity to each.
    """
    text = list(text_lines or [])
    bullets = list(bullet_lines or [])

    if text and bullets:
        t_max, b_max = max_text, max_bullets
        while text or bullets:
            yield (text[:t_max], bullets[:b_max])
            text = text[t_max:]
            bullets = bullets[b_max:]
    elif text:
        while text:
            yield (text[:max_text], [])
            text = text[max_text:]
    elif bullets:
        while bullets:
            yield ([], bullets[:max_bullets])
            bullets = bullets[max_bullets:]
    else:
        yield ([], [])


# -------------------------
# Image helpers
# -------------------------

def _download_if_url(path_or_url: str) -> Optional[str]:
    """
    If path_or_url is a URL and 'requests' is available, download to a temp file and return its path.
    If it's a local path that exists, return it. Otherwise return None.
    """
    if not path_or_url:
        return None
    pu = path_or_url.strip()
    if pu.lower().startswith(("http://", "https://")):
        if not _HAS_REQUESTS:
            return None
        try:
            r = requests.get(pu, timeout=10)
            r.raise_for_status()
            suffix = os.path.splitext(pu)[1] or ".png"
            fd, tmp_path = tempfile.mkstemp(suffix=suffix)
            os.close(fd)
            with open(tmp_path, "wb") as f:
                f.write(r.content)
            return tmp_path
        except Exception:
            return None
    # local file
    return pu if os.path.exists(pu) else None


def _place_image(slide, img_path: str, left, top, max_w, max_h):
    """Place image fit within max_w x max_h, preserve aspect ratio, center inside region."""
    try:
        pic = slide.shapes.add_picture(img_path, left, top)
        # scale down to fit region (never upscale)
        sw = max_w / pic.width
        sh = max_h / pic.height
        scale = min(sw, sh, 1.0)
        pic.width = int(pic.width * scale)
        pic.height = int(pic.height * scale)
        pic.left = left + int((max_w - pic.width) / 2)
        pic.top = top + int((max_h - pic.height) / 2)
        return pic
    except Exception:
        return None


# -------------------------
# Quiz helpers
# -------------------------

def _extract_quiz_choices(block: Dict[str, Any]) -> List[str]:
    """
    Return list [A, B, C, D] if present. Supports:
      - A/B/C/D keys directly
      - choices list on block
    Missing entries become "".
    """
    A = str(block.get("A", "") or "").strip()
    B = str(block.get("B", "") or "").strip()
    C = str(block.get("C", "") or "").strip()
    D = str(block.get("D", "") or "").strip()
    if any([A, B, C, D]):
        return [A, B, C, D]

    ch = block.get("choices")
    if isinstance(ch, list) and ch:
        arr = [str(x or "").strip() for x in ch[:4]]
        while len(arr) < 4:
            arr.append("")
        return arr

    return ["", "", "", ""]


# -------------------------
# Main entry
# -------------------------

def build_storyboard_pptx(
    blocks: List[Dict[str, Any]],
    course_title: str = "Course",
    *,
    font_name: str = "Calibri",
    font_color: str = "#111111",
    bg_color: str = "#FFFFFF",
    max_text_per_slide: int = 6,
    max_bullets_per_slide: int = 6,
) -> io.BytesIO:
    """
    Build a PPTX deck from parsed blocks.

    Slide block shape (all keys optional):
      {
        "type": "slide",
        "title": str,
        "text": [str] or str,
        "bullets": [str] or str,
        "narration": str,      # Notes pane
        "image": str,          # local path or URL
        "alt": str
      }

    Quiz block shape (any subset ok):
      {
        "type": "quiz",
        "title": str,
        "question": str,
        "A": str, "B": str, "C": str, "D": str,
        "choices": [str, str, str, str],   # alternative to A-D
        "answer": "C" or "A,C",            # single or multiple
        "feedback_correct": str,
        "feedback_incorrect": str
      }
    """
    prs = Presentation()
    slide_w = prs.slide_width
    text_rgb = _hex_to_rgb(font_color)
    bg_rgb = _hex_to_rgb(bg_color)

    # --- Title slide ---
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    _apply_slide_bg(title_slide, bg_rgb)

    title_ph = _find_title_placeholder(title_slide)
    if title_ph and title_ph.has_text_frame:
        tf = title_ph.text_frame
        _clear_text_frame(tf)
        p = tf.add_paragraph()
        p.text = course_title
        _apply_font(p, size_pt=44, bold=True, name=font_name, color=text_rgb, bullet=False, level=0)
    else:
        tb = title_slide.shapes.add_textbox(Inches(0.75), Inches(0.7), slide_w - Inches(1.5), Inches(1.1))
        tf = tb.text_frame
        _clear_text_frame(tf)
        _add_para(tf, course_title, size_pt=44, bold=True, name=font_name, color=text_rgb)

    # Optional subtitle
    try:
        sub = title_slide.placeholders[1]
        if sub and sub.has_text_frame:
            _clear_text_frame(sub.text_frame)
            _add_para(sub.text_frame, "Storyboard generated automatically", size_pt=20, bold=False, name=font_name, color=text_rgb)
    except Exception:
        pass

    # --- Content/Quiz slides ---
    for b in blocks:
        if not isinstance(b, dict):
            continue

        btype = b.get("type")

        # ======================
        # QUIZ SLIDE HANDLING
        # ======================
        if btype == "quiz":
            title = (b.get("title") or "Quiz").strip()
            question = str(b.get("question", "") or "")
            answer = str(b.get("answer", "") or "").replace(" ", "")
            feedback_correct = str(b.get("feedback_correct", "") or "")
            feedback_incorrect = str(b.get("feedback_incorrect", "") or "")

            # If multiple letters, show helper text
            if "," in answer:
                question = (question + "\n\n(Select all that apply)").strip()

            # Add slide
            s = prs.slides.add_slide(prs.slide_layouts[1])
            _apply_slide_bg(s, bg_rgb)

            # Title
            tph = _find_title_placeholder(s)
            if tph and tph.has_text_frame:
                tf_title = tph.text_frame
                _clear_text_frame(tf_title)
                _add_para(tf_title, title, size_pt=40, bold=True, name=font_name, color=text_rgb)
            else:
                tb = s.shapes.add_textbox(Inches(0.75), Inches(0.6), slide_w - Inches(1.5), Inches(1.1))
                tf_title = tb.text_frame
                _clear_text_frame(tf_title)
                _add_para(tf_title, title, size_pt=40, bold=True, name=font_name, color=text_rgb)

            # Question + options
            left, top, width, height = _content_bounds(prs)
            tb_q = s.shapes.add_textbox(left, top, width, height)
            tf_q = tb_q.text_frame
            _clear_text_frame(tf_q)

            # Question
            if question:
                _add_para(tf_q, question, size_pt=28, bold=False, name=font_name, color=text_rgb)

            # Options (A-D or choices list)
            A, B, C, D = _extract_quiz_choices(b)
            for letter, text in zip(["A", "B", "C", "D"], [A, B, C, D]):
                if text:
                    _add_bullet(tf_q, f"{letter}: {text}", size_pt=24, level=0, name=font_name, color=text_rgb)

            # Notes with feedback (presenter notes)
            note_text = ""
            if feedback_correct:
                note_text += f"Correct: {feedback_correct}\n"
            if feedback_incorrect:
                note_text += f"Incorrect: {feedback_incorrect}\n"
            if note_text:
                try:
                    notes_tf = s.notes_slide.notes_text_frame
                    _clear_text_frame(notes_tf)
                    for line in note_text.splitlines():
                        _add_para(notes_tf, line, size_pt=14, bold=False, name=font_name, color=text_rgb)
                except Exception:
                    pass
            continue  # next block

        # ======================
        # NORMAL SLIDE HANDLING
        # ======================
        if btype != "slide":
            continue

        title = (b.get("title") or "Slide").strip()
        body_text = _to_lines(b.get("text"))
        bullets = _to_lines(b.get("bullets"))
        narration = (b.get("narration") or "").strip()
        image_ref = (b.get("image") or "").strip()

        # Prepare image (only for first chunk)
        img_path = _download_if_url(image_ref) if image_ref else None

        # Render one or more slides depending on content length
        chunk_idx = 0
        for text_chunk, bullet_chunk in _chunk_content(body_text, bullets, max_text=max_text_per_slide, max_bullets=max_bullets_per_slide):
            chunk_idx += 1
            title_this = title if chunk_idx == 1 else f"{title} (cont.)"

            s = prs.slides.add_slide(prs.slide_layouts[1])
            _apply_slide_bg(s, bg_rgb)

            # Title
            tph = _find_title_placeholder(s)
            if tph and tph.has_text_frame:
                tf_title = tph.text_frame
                _clear_text_frame(tf_title)
                _add_para(tf_title, title_this, size_pt=40, bold=True, name=font_name, color=text_rgb)
            else:
                tb = s.shapes.add_textbox(Inches(0.75), Inches(0.6), slide_w - Inches(1.5), Inches(1.1))
                tf_title = tb.text_frame
                _clear_text_frame(tf_title)
                _add_para(tf_title, title_this, size_pt=40, bold=True, name=font_name, color=text_rgb)

            # Content region (start with safe bounds)
            left, top, width, height = _content_bounds(prs)

            # If we have an image and this is the FIRST chunk, reserve ~40% for the image
            if img_path and chunk_idx == 1:
                img_h = int(height * 0.40)
                _place_image(s, img_path, left, top, width, img_h)
                # move text/bullets below image
                top = top + img_h + Inches(0.2)
                height = height - img_h - Inches(0.2)

            # Decide split between text & bullets
            if text_chunk and bullet_chunk:
                text_height = int(height * 0.55)
                bullets_height = height - text_height - Inches(0.2)
            else:
                text_height = height
                bullets_height = 0

            # Text box (bullets OFF)
            if text_chunk:
                tb_text = s.shapes.add_textbox(left, top, width, text_height)
                tf_text = tb_text.text_frame
                _clear_text_frame(tf_text)
                for line in text_chunk:
                    _add_para(tf_text, line, size_pt=28, bold=False, name=font_name, color=text_rgb)

            # Bullets box (bullets ON)
            if bullet_chunk:
                top_bul = top + (text_height + Inches(0.2) if text_chunk else 0)
                tb_bul = s.shapes.add_textbox(left, top_bul, width, bullets_height or height)
                tf_bul = tb_bul.text_frame
                _clear_text_frame(tf_bul)
                for line in bullet_chunk:
                    _add_bullet(tf_bul, line, size_pt=24, level=0, name=font_name, color=text_rgb)

            # Notes (only on the first chunk, to avoid repetition; change if desired)
            if chunk_idx == 1 and narration:
                try:
                    notes_tf = s.notes_slide.notes_text_frame
                    _clear_text_frame(notes_tf)
                    for line in narration.splitlines():
                        _add_para(notes_tf, line, size_pt=14, bold=False, name=font_name, color=text_rgb)
                except Exception:
                    pass

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
